"""
Sesión 8 de agentes · Generar los entregables (PowerPoint + PDF)

Lee el banco de citas VERIFICADAS (JSON de la sesión 7) y produce:
  1. Un PowerPoint: una cita por diapositiva (frase en inglés + traducción + cita APA).
  2. Un PDF de tarjetas imprimibles: una tarjeta por cita.

Solo usa las citas verificadas. Cada tarjeta lleva su cita en APA con la página.

Requisitos:  uv add python-pptx reportlab
"""

import os
import re
import json

from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from reportlab.lib.pagesizes import A4
from reportlab.lib.units import mm
from reportlab.lib import colors
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle


def _ruta(ruta):
    return os.path.expanduser(ruta.strip())


def limpia(texto):
    """Sustituye comillas y guiones 'tipográficos' por ASCII, para evitar
    problemas de fuentes en el PDF."""
    reemplazos = {"‘": "'", "’": "'", "“": '"', "”": '"',
                  "–": "-", "—": "-", "…": "..."}
    for a, b in reemplazos.items():
        texto = texto.replace(a, b)
    return texto


def cita_corta(ref_apa, pagina):
    """Construye una cita breve tipo (Autor et al., año, p. N) a partir de la
    referencia APA completa. Es una aproximación: revísala si hace falta."""
    anyo = re.search(r"\((\d{4})\)", ref_apa)
    anyo = anyo.group(1) if anyo else "s.f."
    autores = ref_apa.split("(")[0].strip().rstrip(".")
    primer = autores.split(",")[0].strip() if autores else "Autor"
    return f"({primer} et al., {anyo}, p. {pagina})"


# --- 1. POWERPOINT ---
def generar_pptx(banco, salida):
    prs = Presentation()
    prs.slide_width = Inches(13.33)   # formato panorámico 16:9
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]      # diapositiva en blanco

    ref = banco.get("referencia_apa", "")

    # Diapositiva de portada
    s = prs.slides.add_slide(blank)
    caja = s.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.3), Inches(2.5))
    tf = caja.text_frame
    tf.word_wrap = True
    tf.text = "Banco de citas"
    tf.paragraphs[0].font.size = Pt(44)
    tf.paragraphs[0].font.bold = True
    p = tf.add_paragraph()
    p.text = limpia(ref)
    p.font.size = Pt(18)

    # Una diapositiva por cita
    for c in banco.get("citas_verificadas", []):
        s = prs.slides.add_slide(blank)
        caja = s.shapes.add_textbox(Inches(0.8), Inches(0.8), Inches(11.7), Inches(5.9))
        tf = caja.text_frame
        tf.word_wrap = True

        # Cita en inglés (grande)
        tf.text = f'"{limpia(c["cita"])}"'
        tf.paragraphs[0].font.size = Pt(28)
        tf.paragraphs[0].font.bold = True

        # Traducción (mediana, cursiva, gris)
        p = tf.add_paragraph()
        p.text = limpia(c.get("traduccion", ""))
        p.font.size = Pt(18)
        p.font.italic = True
        p.font.color.rgb = RGBColor(0x55, 0x55, 0x55)

        # Cita APA al pie
        pie = s.shapes.add_textbox(Inches(0.8), Inches(6.6), Inches(11.7), Inches(0.6))
        ptf = pie.text_frame
        ptf.text = cita_corta(ref, c.get("pagina", "?"))
        ptf.paragraphs[0].font.size = Pt(14)
        ptf.paragraphs[0].alignment = PP_ALIGN.RIGHT

    prs.save(_ruta(salida))
    return len(banco.get("citas_verificadas", []))


# --- 2. PDF DE TARJETAS ---
def generar_pdf(banco, salida):
    doc = SimpleDocTemplate(_ruta(salida), pagesize=A4,
                            topMargin=15 * mm, bottomMargin=15 * mm,
                            leftMargin=18 * mm, rightMargin=18 * mm)
    estilos = getSampleStyleSheet()
    est_cita = ParagraphStyle("cita", parent=estilos["Normal"],
                              fontName="Helvetica-Bold", fontSize=12, leading=16)
    est_trad = ParagraphStyle("trad", parent=estilos["Normal"],
                              fontName="Helvetica-Oblique", fontSize=10,
                              leading=13, textColor=colors.HexColor("#555555"))
    est_ref = ParagraphStyle("ref", parent=estilos["Normal"],
                             fontName="Helvetica", fontSize=9,
                             textColor=colors.HexColor("#1F4E79"), alignment=2)

    ref = banco.get("referencia_apa", "")
    elementos = [Paragraph("Banco de citas", estilos["Title"]),
                 Paragraph(limpia(ref), est_trad), Spacer(1, 8 * mm)]

    for c in banco.get("citas_verificadas", []):
        contenido = [
            [Paragraph(f'"{limpia(c["cita"])}"', est_cita)],
            [Paragraph(limpia(c.get("traduccion", "")), est_trad)],
            [Paragraph(cita_corta(ref, c.get("pagina", "?")), est_ref)],
        ]
        tarjeta = Table(contenido, colWidths=[doc.width])
        tarjeta.setStyle(TableStyle([
            ("BOX", (0, 0), (-1, -1), 0.75, colors.HexColor("#AACCEE")),
            ("TOPPADDING", (0, 0), (-1, -1), 8),
            ("BOTTOMPADDING", (0, 0), (-1, -1), 8),
            ("LEFTPADDING", (0, 0), (-1, -1), 10),
            ("RIGHTPADDING", (0, 0), (-1, -1), 10),
        ]))
        elementos.append(tarjeta)
        elementos.append(Spacer(1, 6 * mm))

    doc.build(elementos)
    return len(banco.get("citas_verificadas", []))


if __name__ == "__main__":
    carpeta = "~/TESIS doctoral MAS/Publicaciones LEER/0 GRIPE_FLORIN/"
    entrada = carpeta + "Banco de citas - Gripe Florin.json"
    salida_pptx = carpeta + "Citas - Gripe Florin.pptx"
    salida_pdf = carpeta + "Citas - Gripe Florin.pdf"

    with open(_ruta(entrada), "r", encoding="utf-8") as f:
        banco = json.load(f)

    n1 = generar_pptx(banco, salida_pptx)
    print(f"PowerPoint creado con {n1} diapositivas de citas.")
    n2 = generar_pdf(banco, salida_pdf)
    print(f"PDF de tarjetas creado con {n2} tarjetas.")
    print(f"\n==> Archivos en:\n    {carpeta}")
